"""
文件处理工具集
=============
提供文件解析、内容提取、PPT生成、文案润色等工具。

设计决策：
  - 每个工具封装为独立函数，可由多个 Agent 共享
  - 工具返回 Pydantic 模型，确保数据结构一致
  - OCR 作为可选模块，未安装 pytesseract 时优雅降级

可扩展点：
  - 新增文件格式：添加解析函数并在 FILE_PARSERS 注册
  - PPT 生成优化：接入 pptx 模板引擎
  - 润色风格自定义：支持不同 tone (formal/casual/technical)
"""

import logging
from pathlib import Path
from typing import List, Optional, Dict, Any

from src.models.schemas import ParsedFile, FileType

logger = logging.getLogger(__name__)

# ===================== 文件格式检测 =====================

def detect_file_type(filepath: str) -> FileType:
    """根据扩展名检测文件类型"""
    ext = Path(filepath).suffix.lower()
    mapping = {
        ".pdf": FileType.PDF,
        ".docx": FileType.DOCX,
        ".pptx": FileType.PPTX,
        ".xlsx": FileType.XLSX,
        ".csv": FileType.CSV,
        ".png": FileType.PNG,
        ".jpg": FileType.JPG,
        ".jpeg": FileType.JPG,
        ".bmp": FileType.BMP,
        ".tiff": FileType.TIFF,
        ".tif": FileType.TIFF,
        ".webp": FileType.WEBP,
        ".gif": FileType.GIF,
        ".txt": FileType.TXT,
        ".md": FileType.MD,
    }
    return mapping.get(ext, FileType.UNKNOWN)


# ===================== 文件解析器 =====================

def parse_pdf(filepath: str) -> ParsedFile:
    """解析 PDF 文件 — 提取文本、元数据、表格"""
    import PyPDF2

    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.PDF)

    try:
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            parsed.page_count = len(reader.pages)

            # 提取元数据
            if reader.metadata:
                parsed.metadata = {
                    "title": reader.metadata.get("/Title", ""),
                    "author": reader.metadata.get("/Author", ""),
                    "subject": reader.metadata.get("/Subject", ""),
                }

            # 提取文本
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            parsed.raw_text = "\n".join(text_parts)

        logger.info(f"Parsed PDF: {filepath}, pages={parsed.page_count}, chars={len(parsed.raw_text)}")

    except Exception as e:
        logger.error(f"Failed to parse PDF {filepath}: {e}")
        parsed.raw_text = f"[PDF解析失败: {e}]"

    return parsed


def parse_docx(filepath: str) -> ParsedFile:
    """解析 Word 文档"""
    from docx import Document

    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.DOCX)

    try:
        doc = Document(filepath)

        # 提取段落文本
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        parsed.raw_text = "\n".join(paragraphs)

        # 提取表格
        for table in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            parsed.tables.append(rows)

        # 提取元数据
        props = doc.core_properties
        parsed.metadata = {
            "author": str(props.author) if props.author else "",
            "title": str(props.title) if props.title else "",
            "created": str(props.created) if props.created else "",
        }

        logger.info(f"Parsed DOCX: {filepath}, paragraphs={len(paragraphs)}")

    except Exception as e:
        logger.error(f"Failed to parse DOCX {filepath}: {e}")
        parsed.raw_text = f"[DOCX解析失败: {e}]"

    return parsed


def parse_pptx(filepath: str) -> ParsedFile:
    """解析 PPT 文件"""
    from pptx import Presentation

    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.PPTX)

    try:
        prs = Presentation(filepath)
        parsed.page_count = len(prs.slides)

        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text.strip())
                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    parsed.tables.append(rows)
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

        parsed.raw_text = "\n\n".join(text_parts)
        logger.info(f"Parsed PPTX: {filepath}, slides={parsed.page_count}")

    except Exception as e:
        logger.error(f"Failed to parse PPTX {filepath}: {e}")
        parsed.raw_text = f"[PPTX解析失败: {e}]"

    return parsed


def parse_image(filepath: str, ocr_lang: Optional[str] = None, ocr_preprocess: bool = True) -> ParsedFile:
    """
    解析图片文件 — OCR 文字识别

    支持格式: PNG, JPG, JPEG, BMP, TIFF, WEBP, GIF
    OCR 引擎: Tesseract (pytesseract)，自动预处理提高识别率

    可扩展点:
      - 接入更高级的 OCR 引擎 (如 PaddleOCR)
      - 图片分类: 识别是截图/照片/图表
      - 多模态模型直接理解图片内容
    """
    from PIL import Image, ImageEnhance

    file_type = detect_file_type(filepath)
    parsed = ParsedFile(filename=Path(filepath).name, file_type=file_type)

    try:
        img = Image.open(filepath)
        # Ensure RGB mode for consistent processing (GIF, PNG with alpha, etc.)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        parsed.metadata = {"width": img.width, "height": img.height, "format": img.format}

        # OCR 文字识别
        try:
            import pytesseract

            # Pre-processing for better OCR accuracy
            ocr_img = img
            if ocr_preprocess:
                try:
                    # Convert to grayscale
                    ocr_img = img.convert("L")
                    # Enhance contrast
                    enhancer = ImageEnhance.Contrast(ocr_img)
                    ocr_img = enhancer.enhance(2.0)
                except Exception as e:
                    logger.debug(f"OCR pre-processing failed, using original image: {e}")
                    ocr_img = img

            lang = ocr_lang or "chi_sim+eng"
            parsed.raw_text = pytesseract.image_to_string(ocr_img, lang=lang)
            logger.info(f"OCR completed: {filepath}, chars={len(parsed.raw_text)}, lang={lang}")

        except ImportError:
            parsed.raw_text = "[OCR 未安装: pip install pytesseract]"
        except Exception as e:
            parsed.raw_text = f"[OCR 识别失败: {e}]"

        logger.info(f"Parsed Image: {filepath}, size={img.size}, type={file_type.value}")

    except Exception as e:
        logger.error(f"Failed to parse Image {filepath}: {e}")
        parsed.raw_text = f"[图片解析失败: {e}]"

    return parsed


def parse_txt(filepath: str) -> ParsedFile:
    """解析纯文本文件"""
    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.TXT)

    try:
        # 尝试多种编码
        for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
            try:
                with open(filepath, "r", encoding=encoding) as f:
                    parsed.raw_text = f.read()
                parsed.metadata["encoding"] = encoding
                break
            except UnicodeDecodeError:
                continue
        else:
            parsed.raw_text = "[无法识别的文件编码]"

    except Exception as e:
        logger.error(f"Failed to parse TXT {filepath}: {e}")
        parsed.raw_text = f"[文本解析失败: {e}]"

    return parsed


def parse_csv(filepath: str) -> ParsedFile:
    """解析 CSV 文件，输出带表头的表格文本"""
    import csv
    import io

    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.CSV)

    try:
        content = None
        for encoding in ["utf-8", "utf-8-sig", "gbk", "gb2312", "latin-1"]:
            try:
                with open(filepath, "r", encoding=encoding) as f:
                    content = f.read()
                parsed.metadata["encoding"] = encoding
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            parsed.raw_text = "[无法识别的CSV文件编码]"
            return parsed

        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        parsed.metadata["row_count"] = len(rows)
        parsed.metadata["col_count"] = len(rows[0]) if rows else 0

        if not rows:
            parsed.raw_text = "[空CSV文件]"
            return parsed

        col_widths = [0] * len(rows[0])
        for row in rows:
            for i, cell in enumerate(row):
                if i < len(col_widths):
                    col_widths[i] = max(col_widths[i], len(cell))

        lines = []
        for ri, row in enumerate(rows):
            cells = [cell.ljust(col_widths[i]) if i < len(col_widths) else cell
                     for i, cell in enumerate(row)]
            lines.append(" | ".join(cells))
            if ri == 0:
                lines.append("-+-".join("-" * w for w in col_widths))

        parsed.raw_text = "\n".join(lines)

    except Exception as e:
        logger.error(f"Failed to parse CSV {filepath}: {e}")
        parsed.raw_text = f"[CSV解析失败: {e}]"

    return parsed


def parse_xlsx(filepath: str) -> ParsedFile:
    """解析 Excel 文件，提取所有 sheet 的文本内容"""
    import openpyxl

    parsed = ParsedFile(filename=Path(filepath).name, file_type=FileType.XLSX)

    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        parsed.metadata["sheet_count"] = len(wb.sheetnames)
        parsed.metadata["sheets"] = wb.sheetnames

        all_parts = []
        total_rows = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            non_empty = [
                [str(c) if c is not None else "" for c in row]
                for row in rows
                if any(c is not None for c in row)
            ]
            if not non_empty:
                continue

            max_cols = max(len(row) for row in non_empty)
            col_widths = [0] * max_cols
            for row in non_empty:
                for i, cell in enumerate(row):
                    col_widths[i] = max(col_widths[i], len(str(cell)))

            lines = [f"--- Sheet: {sheet_name} ---"]
            for ri, row in enumerate(non_empty):
                cells = [str(cell).ljust(col_widths[i]) if i < len(col_widths) else str(cell)
                         for i, cell in enumerate(row)]
                lines.append(" | ".join(cells))
                if ri == 0:
                    lines.append("-+-".join("-" * w for w in col_widths))

            all_parts.append("\n".join(lines))
            total_rows += len(non_empty)

        parsed.metadata["total_rows"] = total_rows
        parsed.raw_text = "\n\n".join(all_parts) if all_parts else "[空Excel文件]"
        wb.close()

    except Exception as e:
        logger.error(f"Failed to parse XLSX {filepath}: {e}")
        parsed.raw_text = f"[Excel解析失败: {e}]"

    return parsed


# 文件类型 -> 解析器映射 (可扩展点：在此注册新的解析器)
FILE_PARSERS = {
    FileType.PDF: parse_pdf,
    FileType.DOCX: parse_docx,
    FileType.PPTX: parse_pptx,
    FileType.PNG: parse_image,
    FileType.JPG: parse_image,
    FileType.BMP: parse_image,
    FileType.TIFF: parse_image,
    FileType.WEBP: parse_image,
    FileType.GIF: parse_image,
    FileType.TXT: parse_txt,
    FileType.MD: parse_txt,
    FileType.CSV: parse_csv,
    FileType.XLSX: parse_xlsx,
}


def parse_file(filepath: str) -> ParsedFile:
    """
    统一文件解析入口 — 根据文件类型自动选择解析器

    Args:
        filepath: 文件路径

    Returns:
        ParsedFile: 解析后的文件内容
    """
    file_type = detect_file_type(filepath)
    parser = FILE_PARSERS.get(file_type)

    if parser is None:
        logger.warning(f"Unsupported file type: {file_type}")
        return ParsedFile(
            filename=Path(filepath).name,
            file_type=FileType.UNKNOWN,
            raw_text=f"[不支持的文件格式: {filepath}]",
        )

    return parser(filepath)


# ===================== 文件生成工具 =====================

def generate_pptx(
    title: str,
    slides_content: List[Dict[str, Any]],
    output_path: str,
    template_path: Optional[str] = None,
) -> str:
    """
    生成 PPT 文件

    设计决策：
      - 支持模板和从头创建两种模式
      - 内容结构为 [{title: str, bullets: [str], image: Optional[bytes]}, ...]

    Args:
        title: PPT 标题
        slides_content: 每页幻灯片的内容
        output_path: 输出路径
        template_path: 模板文件路径 (可选)

    Returns:
        生成的文件路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    # TODO: 支持模板加载
    prs = Presentation(template_path) if template_path else Presentation()

    for slide_data in slides_content:
        # 使用标题+内容布局
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        # 设置内容
        if slide.placeholders and len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            bullets = slide_data.get("bullets", [])
            text_frame = content.text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0

    # 确保输出目录存在
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    prs.save(output_path)
    logger.info(f"Generated PPTX: {output_path}, slides={len(slides_content)}")
    return output_path


# ===================== 内容处理工具 =====================

def extract_keywords(text: str, top_k: int = 10) -> List[str]:
    """
    从文本中提取关键词

    可扩展点: 使用 TF-IDF 或 LLM 进行智能提取
    """
    # TODO: 实现基于 LLM 的关键词提取
    # 当前返回简单模拟
    words = text.split()
    # 简单过滤短词和重复
    seen = set()
    keywords = []
    for w in words:
        if len(w) > 2 and w not in seen:
            keywords.append(w)
            seen.add(w)
        if len(keywords) >= top_k:
            break
    return keywords
